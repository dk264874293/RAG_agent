"""
工业级文档处理流水线
支持：PDF、Word、Excel、PPT、HTML、Markdown、图片OCR等
"""
import asyncio
import hashlib
from typing import List, Dict, Optional
from pathlib import Path
from ..models.document import Document
from ..extractor.extract_processor import ExtractProcessor
from ..extractor.pdf_extractor import PdfExtractor
from ..extractor.word_extractor import WorkExtractor


class UnsupportedFormatError(Exception):
    """不支持的文件格式异常"""
    pass


class DocumentProcessingPipeline:
    """文档处理流水线"""
    
    def __init__(self, config: Optional[Dict] = None):
        self.config = config or {}
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.html': self._process_html,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
        }
    
    def _detect_format(self, file_path: str) -> str:
        """检测文件格式"""
        return Path(file_path).suffix.lower()
    
    async def process_document(self, file_path: str, metadata: Optional[Dict] = None) -> List[Document]:
        """异步处理文档"""
        # 1. 格式检测与验证
        file_ext = self._detect_format(file_path)
        if file_ext not in self.supported_formats:
            raise UnsupportedFormatError(f"不支持的文件格式: {file_ext}")
        
        # 2. 内容提取
        processor_func = self.supported_formats[file_ext]
        content = await processor_func(file_path)
        
        # 3. 文本清洗和标准化
        cleaned_content = self._clean_content(content)
        
        # 4. 分块策略（将在后续步骤中实现）
        # chunks = await self._chunk_document(cleaned_content, file_ext)
        
        # 5. 元数据增强
        enhanced_metadata = self._enhance_metadata(metadata or {}, file_path, file_ext)
        
        # 创建文档对象
        document = Document(
            page_content=cleaned_content,
            metadata=enhanced_metadata
        )
        
        return [document]
    
    async def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        # 使用现有的PDF提取器
        extractor = PdfExtractor(file_path, '1', '1')
        documents = extractor.extract()
        # 合并所有文档内容
        return "\n\n".join([doc.page_content for doc in documents])
    
    async def _process_docx(self, file_path: str) -> str:
        """处理Word文档"""
        extractor = WorkExtractor(file_path, tenant_id='default', user_id='default')
        documents = extractor.extract()
        return "\n\n".join([doc.page_content for doc in documents])
    
    async def _process_doc(self, file_path: str) -> str:
        """处理旧版Word文档"""
        # 可以调用WordExtractor或使用其他库
        return await self._process_docx(file_path)
    
    async def _process_text(self, file_path: str) -> str:
        """处理纯文本文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    async def _process_markdown(self, file_path: str) -> str:
        """处理Markdown文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    async def _process_html(self, file_path: str) -> str:
        """处理HTML文件"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text()
        except ImportError:
            # 如果没有BeautifulSoup，简单提取文本
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                # 简单去除HTML标签
                import re
                return re.sub(r'<[^>]+>', '', content)
    
    async def _process_pptx(self, file_path: str) -> str:
        """处理PowerPoint文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return "\n\n".join(text_content)
        except ImportError:
            raise UnsupportedFormatError("需要安装python-pptx库来处理PPTX文件")
    
    async def _process_xlsx(self, file_path: str) -> str:
        """处理Excel文件"""
        try:
            import pandas as pd
            df = pd.read_excel(file_path, sheet_name=None)
            text_content = []
            for sheet_name, sheet_df in df.items():
                text_content.append(f"工作表: {sheet_name}\n{sheet_df.to_string()}")
            return "\n\n".join(text_content)
        except ImportError:
            raise UnsupportedFormatError("需要安装pandas和openpyxl库来处理Excel文件")
    
    def _clean_content(self, content: str) -> str:
        """文本清洗和标准化"""
        # 去除多余的空白字符
        import re
        content = re.sub(r'\s+', ' ', content)
        # 去除特殊字符（可根据需要调整）
        content = content.strip()
        return content
    
    def _enhance_metadata(self, metadata: Dict, file_path: str, file_ext: str) -> Dict:
        """元数据增强"""
        enhanced = metadata.copy()
        path_obj = Path(file_path)
        
        # 添加文件信息
        enhanced.update({
            'source': str(file_path),
            'filename': path_obj.name,
            'file_type': file_ext,
            'file_size': path_obj.stat().st_size if path_obj.exists() else 0,
            'file_hash': self._calculate_file_hash(file_path) if path_obj.exists() else None,
        })
        
        return enhanced
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
